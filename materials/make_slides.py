# -*- coding: utf-8 -*-
"""
make_slides.py
직업계고 1학년 「국어 x 디자인」 읽기 독해 수업용 슬라이드(.pptx, 11장, 16:9) 생성 스크립트.

명세(SSOT): 03_PPT.md / 방법론: .claude/skills/pptx-slides/SKILL.md
공통 맥락: _workspace/00_shared-brief.md (AI 윤리 3원칙: 개인정보 금지 / 원문 대조 / 보조 도구로만)

실행:  python materials\\make_slides.py
출력:  materials/output/수업PPT_디자인직무읽기.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------------------------------------------------------------- 색/톤 (네이비·블루 단일 톤)
NAVY  = RGBColor(0x1F, 0x3A, 0x5F)
BLUE  = RGBColor(0x2D, 0x6C, 0xDF)
INK   = RGBColor(0x22, 0x2A, 0x33)
GRAY  = RGBColor(0x6B, 0x73, 0x80)
LIGHT = RGBColor(0xEC, 0xF1, 0xFA)   # 연한 블루 배경
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT  = "맑은 고딕"

# ---------------------------------------------------------------- 캔버스 16:9
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = 13.333, 7.5


def set_font(run, size, color=INK, bold=False, name=FONT):
    """모든 run에 한글 폰트를 라틴·eastAsia 양쪽으로 지정해 깨짐 방지."""
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name
    rPr = run.font._rPr
    # East Asian 폰트 보장
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', name)
    # 라틴 폰트도 명시
    latin = rPr.find(qn('a:latin'))
    if latin is None:
        latin = rPr.makeelement(qn('a:latin'), {})
        rPr.insert(0, latin)
    latin.set('typeface', name)


def add_slide():
    return prs.slides.add_slide(BLANK)


def bg(slide, color=WHITE):
    """슬라이드 전체 배경색."""
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                               prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background()
    r.shadow.inherit = False
    slide.shapes._spTree.remove(r._element)
    slide.shapes._spTree.insert(2, r._element)  # 맨 뒤로
    return r


def rect(slide, l, t, w, h, color, line_color=None):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t),
                               Inches(w), Inches(h))
    r.fill.solid(); r.fill.fore_color.rgb = color
    if line_color is None:
        r.line.fill.background()
    else:
        r.line.color.rgb = line_color; r.line.width = Pt(1)
    r.shadow.inherit = False
    return r


def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    return tf


def para(tf, text, size, color=INK, bold=False, align=PP_ALIGN.LEFT,
         space_after=10, space_before=0, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    run = p.add_run(); run.text = text
    set_font(run, size, color, bold)
    return p


def title_band(slide, title, sub=None):
    """상단 NAVY 컬러 바 + 제목(공통 본문 슬라이드 헤더)."""
    rect(slide, 0, 0, SW, 1.35, NAVY)
    rect(slide, 0, 1.35, SW, 0.08, BLUE)      # 얇은 블루 강조선
    tf = textbox(slide, 0.7, 0.18, SW - 1.4, 1.0, MSO_ANCHOR.MIDDLE)
    para(tf, title, 34, WHITE, bold=True, first=True, space_after=2)
    if sub:
        para(tf, sub, 16, RGBColor(0xC7, 0xD6, 0xF0), space_after=0)


def badge(slide, label):
    """활동 번호 배지(둥근 사각) — 좌상단 단계 표시."""
    b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(0.7), Inches(1.75), Inches(2.5), Inches(0.62))
    b.fill.solid(); b.fill.fore_color.rgb = BLUE
    b.line.fill.background(); b.shadow.inherit = False
    tf = b.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, label, 20, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True, space_after=0)


def bullets(slide, items, top=2.7, left=0.9, width=11.5, size=24, gap=16):
    """글머리 본문(3~5개, 큰 글씨, 넉넉한 간격)."""
    tf = textbox(slide, left, top, width, SH - top - 0.5)
    for i, (txt, *opt) in enumerate(items):
        sub = opt[0] if opt else False
        if sub:
            p = para(tf, "   - " + txt, size - 4, GRAY, first=(i == 0), space_after=gap - 4)
        else:
            p = para(tf, "•  " + txt, size, INK, bold=False, first=(i == 0), space_after=gap)
    return tf


# ================================================================ 슬라이드 1: 표지
def slide_01():
    s = add_slide(); bg(s, NAVY)
    rect(s, 0, 0, SW, SH, NAVY)
    rect(s, 0, 5.05, SW, 0.06, BLUE)
    # 상단 라벨
    tf = textbox(s, 1.0, 1.5, SW - 2.0, 0.6)
    para(tf, "직업계고 1학년 · 국어 × 디자인", 20, RGBColor(0x9F, 0xBE, 0xEC), bold=True, first=True)
    # 메인 타이틀
    tf = textbox(s, 1.0, 2.35, SW - 2.0, 2.4, MSO_ANCHOR.TOP)
    para(tf, "AI 독해 도우미로", 50, WHITE, bold=True, first=True, space_after=6)
    para(tf, "디자인 직무 글 읽기", 50, WHITE, bold=True, space_after=0)
    # 부제
    tf = textbox(s, 1.0, 5.3, SW - 2.0, 1.2)
    para(tf, "잘 만들려면, 먼저 잘 읽어야 합니다", 22, RGBColor(0xC7, 0xD6, 0xF0), first=True, space_after=4)
    para(tf, "[10공국1-02-03] 진로·관심 분야 글을 읽고 읽기 과정을 점검·조정하기", 15, RGBColor(0x8F, 0xA8, 0xCF), space_after=0)


# ================================================================ 슬라이드 2: 동기유발
def slide_02():
    s = add_slide(); bg(s)
    title_band(s, "디자인 현장의 '읽기'", "디자이너는 생각보다 글을 많이 읽습니다")
    bullets(s, [
        ("트렌드 리포트 — 올해 유행하는 색·스타일 파악",),
        ("클라이언트 브리프 — 고객이 원하는 것 이해",),
        ("작업 지시서 — 무엇을, 어떻게 만들지 확인",),
        ("소재·공정 설명서 — 재료와 만드는 과정 익히기",),
    ], top=2.55, size=25, gap=18)
    # 핵심 메시지 박스
    rect(s, 0.9, 6.05, 11.5, 0.95, LIGHT)
    rect(s, 0.9, 6.05, 0.12, 0.95, BLUE)
    tf = textbox(s, 1.2, 6.05, 11.0, 0.95, MSO_ANCHOR.MIDDLE)
    para(tf, "→ \"잘 만들려면 잘 읽어야 한다\"", 24, NAVY, bold=True, first=True, space_after=0)


# ================================================================ 슬라이드 3: 오늘의 목표
def slide_03():
    s = add_slide(); bg(s)
    title_band(s, "오늘의 목표", "이 수업이 끝나면 이렇게 할 수 있어요")
    bullets(s, [
        ("디자인 직무 글을 읽고 핵심 정보를 찾을 수 있다",),
        ("글쓴이가 말하려는 의도를 파악할 수 있다",),
        ("AI 독해 도우미로 내 읽기 과정을 점검할 수 있다",),
    ], top=2.85, size=27, gap=26)


# ================================================================ 슬라이드 4: AI 도우미 소개 (+QR)
def slide_04():
    s = add_slide(); bg(s)
    title_band(s, "AI 독해 도우미 소개", "읽기를 도와주는 4가지 기능")
    # 왼쪽: 기능 + 사용법
    tf = textbox(s, 0.9, 2.45, 7.7, 4.6)
    para(tf, "이런 걸 도와줘요", 22, BLUE, bold=True, first=True, space_after=10)
    for txt in ["① 어휘 풀이 — 모르는 낱말 뜻 알려주기",
                "② 문단 요약 — 긴 문단을 짧게 정리",
                "③ 핵심어 — 중요한 단어 짚어주기",
                "④ 이해 점검 질문 — 내가 잘 읽었는지 확인"]:
        para(tf, txt, 21, INK, space_after=8)
    para(tf, "사용법 3단계", 22, BLUE, bold=True, space_before=10, space_after=10)
    for txt in ["1) 웹앱 접속 → 2) 글 붙여넣기 → 3) 기능 선택"]:
        para(tf, txt, 21, INK, space_after=0)

    # 오른쪽: QR / 웹앱 주소 placeholder (점선 회색 사각형 + 안내)
    qr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.2), Inches(2.55),
                            Inches(3.2), Inches(3.2))
    qr.fill.solid(); qr.fill.fore_color.rgb = RGBColor(0xF2, 0xF4, 0xF7)
    qr.line.color.rgb = GRAY; qr.line.width = Pt(1.5)
    qr.line.dash_style = 2  # 점선
    qr.shadow.inherit = False
    tf = qr.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, "[ QR 코드 자리 ]", 18, GRAY, bold=True, align=PP_ALIGN.CENTER, first=True, space_after=6)
    para(tf, "여기에 웹앱 주소·QR를\n넣어 주세요", 15, GRAY, align=PP_ALIGN.CENTER, space_after=0)
    # 주소 입력 안내선
    tf = textbox(s, 9.2, 5.95, 3.2, 0.7)
    para(tf, "웹앱 주소:", 15, NAVY, bold=True, first=True, space_after=2)
    para(tf, "____________________", 15, GRAY, space_after=0)


# ================================================================ 슬라이드 5: AI 사용 주의점
def slide_05():
    s = add_slide(); bg(s)
    title_band(s, "AI 사용 주의점", "AI는 똑똑한 도우미일 뿐, 4가지를 꼭 지켜요")
    cards = [
        ("①", "개인정보 입력 금지", "이름·연락처 등 내 정보를\n절대 넣지 않기"),
        ("②", "AI는 틀릴 수 있다", "답을 그대로 믿지 말고\n원문과 꼭 대조하기"),
        ("③", "'보조 도구'로만", "AI는 도움일 뿐,\n읽는 사람은 바로 나"),
        ("④", "베끼지 말기", "AI 답을 복사하지 말고\n내 말로 다시 정리"),
    ]
    x0, y0, cw, ch, gapx, gapy = 0.9, 2.55, 5.6, 1.95, 0.33, 0.3
    for i, (num, head, body) in enumerate(cards):
        col, row = i % 2, i // 2
        l = x0 + col * (cw + gapx)
        t = y0 + row * (ch + gapy)
        c = rect(s, l, t, cw, ch, LIGHT)
        rect(s, l, t, 0.12, ch, BLUE)
        tf = textbox(s, l + 0.35, t + 0.18, cw - 0.55, ch - 0.3, MSO_ANCHOR.MIDDLE)
        para(tf, f"{num} {head}", 23, NAVY, bold=True, first=True, space_after=4)
        para(tf, body, 17, INK, space_after=0)


# ================================================================ 슬라이드 6~9: 활동
def activity_slide(label, title, sub, items, note=None):
    s = add_slide(); bg(s)
    title_band(s, title, sub)
    badge(s, label)
    bullets(s, items, top=2.7, size=24, gap=18)
    if note:
        rect(s, 0.9, 6.2, 11.5, 0.85, LIGHT)
        rect(s, 0.9, 6.2, 0.12, 0.85, BLUE)
        tf = textbox(s, 1.2, 6.2, 11.0, 0.85, MSO_ANCHOR.MIDDLE)
        para(tf, note, 18, NAVY, bold=True, first=True, space_after=0)
    return s


def slide_06():
    activity_slide("활동 1", "스스로 읽기", "AI 도움 없이, 먼저 내 힘으로",
        [("도움 없이 글을 끝까지 천천히 읽기",),
         ("모르는 낱말·이해 안 되는 부분에 표시하기",),
         ("이 글이 무엇에 관한 글인지 한 번 생각해 보기",)],
        note="틀려도 괜찮아요. 먼저 스스로 부딪혀 보는 게 중요해요.")


def slide_07():
    activity_slide("활동 2", "AI 도우미로 비계 받기", "막힌 부분을 AI 도움으로 넘기기",
        [("어휘 확인 — 표시한 낱말의 뜻 찾기",),
         ("문단 요약·핵심어 — 긴 문단을 짧게, 중요한 말 짚기",),
         ("점검 질문 풀기 — 내가 잘 읽었는지 확인",),
         ("AI 답은 원문과 대조하기 (틀릴 수 있어요)", True)])


def slide_08():
    activity_slide("활동 3", "핵심 정리", "읽은 내용을 내 말로 정리하기",
        [("중심 내용 — 이 글에서 가장 중요한 것",),
         ("글쓴이 의도 — 무엇을 말하려고 썼을까",),
         ("디자인 직무 적용 — 내 작업에 어떻게 쓸까",)],
        note="AI 답을 베끼지 말고, 꼭 '내 말'로 정리해요.")


def slide_09():
    activity_slide("활동 4", "짝·모둠 점검", "친구와 비교하며 이해 넓히기",
        [("짝·모둠과 정리한 내용 서로 비교하기",),
         ("다르게 이해한 부분을 찾아 토의하기",),
         ("왜 다르게 읽었는지 이유 이야기하기",)])


# ================================================================ 슬라이드 10: 정리
def slide_10():
    s = add_slide(); bg(s)
    title_band(s, "정리", "오늘 읽기를 돌아봐요")
    bullets(s, [
        ("오늘 읽은 글을 '한 문장'으로 요약해 발표하기",),
        ("소감 나누기 — AI 도우미가 준 도움은 무엇이었나",),
        ("소감 나누기 — AI 도우미의 한계는 무엇이었나",),
    ], top=2.85, size=25, gap=24)
    rect(s, 0.9, 6.15, 11.5, 0.9, LIGHT)
    rect(s, 0.9, 6.15, 0.12, 0.9, BLUE)
    tf = textbox(s, 1.2, 6.15, 11.0, 0.9, MSO_ANCHOR.MIDDLE)
    para(tf, "AI는 도움을 주지만, 끝까지 읽고 판단하는 건 '나'예요.", 19, NAVY, bold=True, first=True, space_after=0)


# ================================================================ 슬라이드 11: 다음 시간
def slide_11():
    s = add_slide(); bg(s, NAVY)
    rect(s, 0, 0, SW, SH, NAVY)
    tf = textbox(s, 1.0, 1.6, SW - 2.0, 0.7)
    para(tf, "다음 시간", 22, RGBColor(0x9F, 0xBE, 0xEC), bold=True, first=True, space_after=0)
    tf = textbox(s, 1.0, 2.5, SW - 2.0, 2.0)
    para(tf, "읽은 자료로", 40, WHITE, bold=True, first=True, space_after=4)
    para(tf, "디자인 트렌드 분석하기", 40, WHITE, bold=True, space_after=0)
    tf = textbox(s, 1.0, 4.7, SW - 2.0, 1.8)
    para(tf, "• 오늘 읽은 글에서 디자인 트렌드 찾아 분석", 22, RGBColor(0xC7, 0xD6, 0xF0), first=True, space_after=10)
    para(tf, "• 직업 분야별 콘셉트 기획하고 발표하기", 22, RGBColor(0xC7, 0xD6, 0xF0), space_after=0)


# ================================================================ 빌드
def build():
    out_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "output")
    os.makedirs(out_dir, exist_ok=True)
    for fn in [slide_01, slide_02, slide_03, slide_04, slide_05,
               slide_06, slide_07, slide_08, slide_09, slide_10, slide_11]:
        fn()
    out_path = os.path.join(out_dir, "수업PPT_디자인직무읽기.pptx")
    prs.save(out_path)
    print("저장 완료:", out_path)
    print("슬라이드 수:", len(prs.slides._sldIdLst))


if __name__ == "__main__":
    build()
