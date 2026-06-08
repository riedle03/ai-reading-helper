# -*- coding: utf-8 -*-
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
from docx import Document
from docx.oxml.ns import qn
from pptx import Presentation
from pptx.util import Emu

DOCX = r"materials/output/활동지_디자인직무읽기.docx"
PPTX = r"materials/output/수업PPT_디자인직무읽기.pptx"

def detect_broken(s):
    bad = [c for c in s if c in ("□","�")]
    return len(bad)

print("="*70)
print("DOCX EXTRACTION:", DOCX)
print("="*70)
try:
    d = Document(DOCX)
    paras = [p.text for p in d.paragraphs if p.text.strip()]
    print("PARA_COUNT:", len(paras))
    print("TABLE_COUNT:", len(d.tables))
    full = "\n".join(p.text for p in d.paragraphs)
    # tables text
    tbl_texts = []
    for ti, t in enumerate(d.tables):
        cells = [c.text for row in t.rows for c in row.cells]
        tbl_texts.append((ti, len(t.rows), len(t.columns), cells))
    alltext = full + "\n" + "\n".join(" ".join(c) for (_,_,_,c) in tbl_texts)
    print("BROKEN_CHARS_TOTAL:", detect_broken(alltext))
    print("---ALL PARAGRAPHS---")
    for p in paras:
        print("P|", p)
    print("---TABLES---")
    for ti,(idx,r,col,cells) in enumerate(tbl_texts):
        print(f"TABLE[{idx}] rows={r} cols={col}")
        for c in cells:
            if c.strip():
                print("  C|", c.replace("\n"," / "))
    # font sample check: scan runs for eastAsia rFonts
    print("---FONT SAMPLE (eastAsia)---")
    seen = set()
    cnt = 0
    for p in d.paragraphs:
        for run in p.runs:
            rPr = run._element.find(qn('w:rPr'))
            if rPr is not None:
                rFonts = rPr.find(qn('w:rFonts'))
                if rFonts is not None:
                    ea = rFonts.get(qn('w:eastAsia'))
                    if ea and ea not in seen:
                        seen.add(ea); cnt+=1
                        print("  eastAsia=", ea, "| sample:", run.text[:30])
        if cnt>=8: break
    if not seen:
        print("  (no eastAsia rFonts found on paragraph runs - checking table runs)")
        for t in d.tables:
            for row in t.rows:
                for cell in row.cells:
                    for p in cell.paragraphs:
                        for run in p.runs:
                            rPr = run._element.find(qn('w:rPr'))
                            if rPr is not None:
                                rFonts = rPr.find(qn('w:rFonts'))
                                if rFonts is not None:
                                    ea = rFonts.get(qn('w:eastAsia'))
                                    if ea and ea not in seen:
                                        seen.add(ea)
                                        print("  (table) eastAsia=", ea)
    print("EASTASIA_FONTS_FOUND:", sorted(seen))
except Exception as e:
    print("DOCX_ERROR:", repr(e))

print()
print("="*70)
print("PPTX EXTRACTION:", PPTX)
print("="*70)
try:
    p = Presentation(PPTX)
    print("SLIDE_COUNT:", len(p.slides))
    print("SLIDE_WIDTH_EMU:", p.slide_width, "HEIGHT_EMU:", p.slide_height)
    ratio = p.slide_width / p.slide_height
    print("ASPECT_RATIO:", round(ratio,4), "(16:9 = 1.7778)")
    allp = []
    for i, s in enumerate(p.slides, 1):
        texts = []
        for sh in s.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                texts.append(sh.text_frame.text)
        joined = " || ".join(texts)
        allp.append(joined)
        print(f"--- SLIDE {i} ---")
        print(joined)
    print("BROKEN_CHARS_TOTAL:", detect_broken("\n".join(allp)))
    # font sample for ppt
    print("---PPT FONT SAMPLE (a:ea / latin)---")
    found = set()
    for s in p.slides:
        for sh in s.shapes:
            if not sh.has_text_frame: continue
            for para in sh.text_frame.paragraphs:
                for run in para.runs:
                    rPr = run._r.find(qn('a:rPr'))
                    if rPr is not None:
                        ea = rPr.find(qn('a:ea'))
                        latin = rPr.find(qn('a:latin'))
                        if ea is not None and ea.get('typeface'):
                            found.add("ea:"+ea.get('typeface'))
                        if latin is not None and latin.get('typeface'):
                            found.add("latin:"+latin.get('typeface'))
    print("PPT_FONTS_FOUND:", sorted(found))
except Exception as e:
    print("PPTX_ERROR:", repr(e))
